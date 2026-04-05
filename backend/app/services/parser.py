import pdfplumber
from pptx import Presentation
import io

def extract_text_from_pdf(content: bytes) -> str:
    text = ""
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for page in pdf.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
    return text

def extract_text_from_pptx(content: bytes) -> str:
    text = ""
    prs = Presentation(io.BytesIO(content))
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_txt(content: bytes) -> str:
    return content.decode("utf-8", errors="ignore")

def extract_text(filename: str, content: bytes, content_type: str) -> str:
    if filename.lower().endswith(".pdf") or "pdf" in content_type:
        return extract_text_from_pdf(content)
    elif filename.lower().endswith(".pptx") or "presentation" in content_type:
        return extract_text_from_pptx(content)
    else:
        # Fallback to plain text
        return extract_text_from_txt(content)
